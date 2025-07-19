# Copyright (c) Meta Platforms, Inc. and affiliates.
# All rights reserved.
#
# This source code is licensed under the terms described in the LICENSE file in
# the root directory of this source tree.

import io
from typing import List, Dict, Any

import fitz  # PyMuPDF
from PIL import Image


import os
import docx
from pptx import Presentation

class MultimodalParser:
    def parse(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Parses a file, extracting text and images.

        Args:
            file_path (str): The path to the file.

        Returns:
            List[Dict[str, Any]]: A list of dictionaries, where each dictionary
                                 represents a page with its text and a single image.
        """
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            return self._parse_pdf(file_path)
        elif ext == ".docx":
            return self._parse_docx(file_path)
        elif ext == ".pptx":
            return self._parse_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file extension for multimodal parsing: {ext}")

    def _parse_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        doc = fitz.open(file_path)
        data = []

        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text()
            image_list = page.get_images(full=True)

            if not image_list:
                data.append({"text": text, "image": None})

            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                data.append({"text": text, "image": image_bytes})

        return data

    def _parse_docx(self, file_path: str) -> List[Dict[str, Any]]:
        doc = docx.Document(file_path)
        data = []
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"

        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image_bytes = rel.target_part.blob
                data.append({"text": text, "image": image_bytes})

        if not data:
            data.append({"text": text, "image": None})

        return data

    def _parse_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        prs = Presentation(file_path)
        data = []
        for slide in prs.slides:
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    image_bytes = shape.image.blob
                    data.append({"text": text, "image": image_bytes})

        if not data:
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            data.append({"text": text, "image": None})

        return data
