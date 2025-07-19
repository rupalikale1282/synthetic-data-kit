"""Integration tests for the multimodal workflow."""

import os
import tempfile
from unittest.mock import MagicMock, patch

import pytest
import fitz  # PyMuPDF
from PIL import Image

from synthetic_data_kit.core import ingest, create


def _create_dummy_pdf(pdf_path, text, image_path):
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 50), text)
    rect = fitz.Rect(50, 100, 150, 200)
    page.insert_image(rect, filename=image_path)
    doc.save(pdf_path)
    doc.close()


def _create_dummy_docx(docx_path, text, image_path):
    from docx import Document
    document = Document()
    document.add_paragraph(text)
    document.add_picture(image_path)
    document.save(docx_path)


def _create_dummy_pptx(pptx_path, text, image_path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    left = top = Inches(1)
    slide.shapes.add_textbox(left, top, Inches(8), Inches(1)).text = text
    slide.shapes.add_picture(image_path, left, top + Inches(2))
    prs.save(pptx_path)


@pytest.mark.integration
@pytest.mark.parametrize("file_type, create_dummy_file", [
    ("pdf", _create_dummy_pdf),
    ("docx", _create_dummy_docx),
    ("pptx", _create_dummy_pptx),
])
def test_multimodal_flow(patch_config, test_env, file_type, create_dummy_file):
    """Test the full multimodal workflow from ingestion to creation."""
    with tempfile.NamedTemporaryFile(suffix=f".{file_type}", delete=False) as dummy_file, \
         tempfile.NamedTemporaryFile(suffix=".png", delete=False) as img_file:
        # Create a dummy image
        img = Image.new('RGB', (100, 100), color = 'red')
        img.save(img_file.name)

        # Create a dummy file
        create_dummy_file(dummy_file.name, "This is a test.", img_file.name)

        ingest_output_dir = tempfile.mkdtemp()
        create_output_dir = tempfile.mkdtemp()

        try:
            # Ingest the file with multimodal flag
            lance_path = ingest.process_file(
                dummy_file.name,
                output_dir=ingest_output_dir,
                multimodal=True
            )
            assert os.path.exists(lance_path)
            assert lance_path.endswith(".lance")

            # Mock LLMClient and VQAGenerator for the create step
            with patch("synthetic_data_kit.core.create.LLMClient") as mock_llm_client_class, \
                 patch("synthetic_data_kit.core.create.VQAGenerator") as mock_vqa_gen_class:
                mock_llm_client = MagicMock()
                mock_llm_client_class.return_value = mock_llm_client

                mock_generator = MagicMock()
                output_path = os.path.join(create_output_dir, "data.parquet")
                mock_generator.process_dataset.return_value = output_path
                mock_vqa_gen_class.return_value = mock_generator

                # Create a dummy parquet file to be returned by the mock
                os.makedirs(create_output_dir, exist_ok=True)
                with open(output_path, "w") as f:
                    f.write("dummy content")

                # Create VQA pairs from the Lance dataset
                parquet_path = create.process_file(
                    lance_path,
                    output_dir=create_output_dir,
                    content_type="vqa"
                )
                assert os.path.exists(parquet_path)
                assert parquet_path.endswith(".parquet")

                mock_generator.process_dataset.assert_called_once()

        finally:
            # Clean up temporary files and directories
            if os.path.exists(dummy_file.name):
                os.unlink(dummy_file.name)
            if os.path.exists(img_file.name):
                os.unlink(img_file.name)
            if os.path.exists(ingest_output_dir):
                import shutil
                shutil.rmtree(ingest_output_dir)
            if os.path.exists(create_output_dir):
                import shutil
                shutil.rmtree(create_output_dir)
